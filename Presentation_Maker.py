from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches
import os, sys, subprocess

class Presentation_Maker:
    def __init__(self):
        self.prs = Presentation()

    def add_slide(self, title, text_list):
        self.slide_layout = self.prs.slide_layouts[1]
        self.slide = self.prs.slides.add_slide(self.slide_layout)
        self.shapes = self.slide.shapes

        self.title_shape = self.shapes.title
        self.body_shape = self.shapes.placeholders[1]

        self.title_shape.text = title
        self.tf = self.body_shape.text_frame

        for x in range(0, len(text_list)):
            self.p = self.tf.add_paragraph()
            self.p.text = text_list[x]
            self.p.level = x
            self.p.alignment = PP_ALIGN.RIGHT
            self.p.font.language_id = MSO_LANGUAGE_ID.HEBREW

    def add_image_slide(self, image_path):
        self.blank_slide = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(self.blank_slide)
        self.image = self.slide.shapes.add_picture(image_path, Inches(
            0.5), Inches(1.75), width=Inches(9), height=Inches(5))

    def save_prs(self, name):
        self.prs.save(name + ".pptx")
        if sys.platform == "win32": # windows
            os.startfile(name + ".pptx")
        else:
            subprocess.call(['xdg-open', name + ".pptx"]) # linux


if __name__ == '__main__':
    presentatia = Presentation_Maker()
    presentatia.add_slide('test', ['שלום עולם, this is a בדיקה'])
    presentatia.save_prs("testN332")
